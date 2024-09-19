import numpy as np
from matplotlib.ticker import FormatStrFormatter
from misc_classes.EndOfLifeProcessing import EndOfLifeProcessing
import matplotlib.pyplot as plt
import os
import pickle
from pd2ppt import df_to_powerpoint, df_to_table
from pptx import Presentation
from scipy.stats import norm
from check_current_os import get_base_path_batt_lab_data
import time
fontProperties = {'family': 'sans-serif',
                  'sans-serif': ['Helvetica'],
                  'weight': 'normal',
                  'size': 12}


def check_or_make_dir(path):
    if not os.path.isdir(path):
        os.mkdir(path)
        print(f'Directory generated for path {path}')


plt.rcParams.update({
    "text.usetex": True,
    "font.family": "Helvetica"
})
BASE_BATLAB_DATA_DIR = get_base_path_batt_lab_data()
test_cases = ["Test1_1", "Test1_2", "Test2_1", "Test2_2"]
TEST_NAMES = ["Test_A", "Test_B", "Test_C", "Test_D"]
name_dict = dict(zip(test_cases, TEST_NAMES))
CASE_NAMES = ['reduce', 'increase']
CASE_SCALES = [1, 1.5, 2]
tic = time.time()
SYNTHETIC_DATA_BASE_DIR = os.path.join(BASE_BATLAB_DATA_DIR, 'stat_test/synthetic_data_base_dir')
for CASE_NAME in CASE_NAMES:
    for CASE_SCALE in CASE_SCALES:
        for tc in TEST_NAMES:
            data_file = os.path.join(SYNTHETIC_DATA_BASE_DIR, f"t_to_eol_sc{CASE_SCALE}/t_to_eol_{tc}_{CASE_NAME}.pkl")
            with open(data_file, 'rb') as h:
                dataset = pickle.load(h)
            eol_pp = EndOfLifeProcessing(dataset)

            # OUTPUT FIGURES FOR HISTOGRAMS AND NORMAL FITS FOR ALL EOL ESTIMATIONS
            update_figs = 1
            shade_area = 1
            fig_base_dir = fr"Z:\StatisticalTest\figures_from_synthetic_data_{CASE_NAME}\scaling_{CASE_SCALE}_new"
            check_or_make_dir(fig_base_dir)
            fig_dir = os.path.join(fig_base_dir, f"{tc}")
            if update_figs:
                print('Updating figures')
                check_or_make_dir(fig_dir)
                for k in eol_pp.dataset.keys():
                    fig, ax = plt.subplots(1, 1)
                    fig_name = os.path.join(fig_dir, f'norm_fit_and_histogram_{k}__{CASE_NAME}.pdf')
                    ax = eol_pp._plot_histogram(ax, normal_bool=True, n_cases=k, col_oe=1)
                    ax = eol_pp._plot_norm_distribution(ax, n_cases=k)
                    ax.set_xticks(ax.get_xticks(), labels=[f'{x:.0f}' for x in ax.get_xticks()])
                    ax.set_yticks(ax.get_yticks(), labels=[f'{y:.3f}' for y in ax.get_yticks()])
                    ax.text(0.05, 0.9, f'{tc.replace("_", " ")}: {k.split("_")[0]} cells',
                            fontsize=12,
                            horizontalalignment='left',
                            verticalalignment='center',
                            transform=ax.transAxes)
                    ax.set_xlabel('FCE to EOL [-]', fontsize=14)
                    ax.set_ylabel('Frequency', fontsize=14)
                    ax.lines[0].set_label(None)
                    ax.legend(loc='upper right')
                    ax.yaxis.set_major_formatter(FormatStrFormatter('%.3f'))
                    if shade_area:
                        mu, sig = eol_pp.norm_dist_params[k]
                        x = np.linspace(mu - 3*sig, mu + 3*sig, 1000)
                        y = norm.pdf(x, loc=mu, scale=sig)
                        rsk = 1 - norm.cdf(1.05 * eol_pp.ref_value, loc=mu, scale=sig)
                        ax.text(0.7, 0.5, f'$P(OE)={100*rsk:.2f}\%$',
                                fontsize=12,
                                horizontalalignment='left',
                                verticalalignment='center',
                                transform=ax.transAxes)
                        ax.fill_between(x[x > 1.05*eol_pp.ref_value], y[x > 1.05*eol_pp.ref_value], 0,
                                        color='orange', alpha=0.5)
                        fig_name = os.path.join(fig_dir, f'norm_fit_and_histogram_{k}_{CASE_NAME}_with_shade_with_cdf.png')
                    fig.savefig(fig_name, dpi=300)
                    fig.savefig(fig_name.replace('png', 'pdf'))
                plt.close('all')

            # OUTPUT RISK BASED ON NORMAL DISTRIBUTION TO TABLE IN POWERPOINT AND LATEX
            output_ppt_table = 1
            LATEX_OUTPUT_DIR = os.path.join(fig_base_dir, 'latex_tables')
            check_or_make_dir(LATEX_OUTPUT_DIR)
            PPT_OUTPUT_DIR = os.path.join(fig_base_dir, 'ppt_tables')
            check_or_make_dir(PPT_OUTPUT_DIR)
            if output_ppt_table:
                df = eol_pp.df_oe_risk_norm
                df.fillna(0, inplace=True)
                df.columns = [c.replace('_test', '') for c in df.columns]
                df = df * 100
                df.insert(0, 'OE', df.index)
                m, n = df.shape
                col_format = ['.3', *['.2' for k in range(1, n)]]
                with open(os.path.join(LATEX_OUTPUT_DIR, f'risk_table_{tc}_{CASE_NAME}.tex'), 'w') as tf:
                    tf.write(df.loc[:, :'7'].to_latex(float_format="%.1f"))
                df_to_powerpoint(os.path.join(PPT_OUTPUT_DIR, f'risk_table_{tc}_{CASE_NAME}.pptx'),
                                 df.iloc[:, :-1],
                                 col_formatters=col_format,
                                 width=19,
                                 height=6,
                                 left=3,
                                 top=3,
                                 name='risk_table')
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                pr_shp = df_to_table(slide, df, col_formatters=col_format)
toc = time.time()
print(f'Elapsed time is {toc - tic:.2f} seconds')
