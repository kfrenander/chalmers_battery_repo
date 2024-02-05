import numpy as np
import os
import re
from misc_classes.PopulationParameterEstimation import PopulationParameterEstimator as pp_itm
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.ticker import FormatStrFormatter
from pd2ppt import df_to_powerpoint, df_to_table
from pptx import Presentation


def _clean_xticks(ax):
    lbls = _clean_labels(ax.get_xticklabels())
    ax.set_xticks(ax.get_xticks())
    ax.set_xticklabels(lbls)
    return ax


def _clean_labels(tck_lbls):
    return [re.sub("[^0-9]", "", itm.get_text()) for itm in tck_lbls]


save_fig = 1
base_dir = r"Z:\StatisticalTest\figures_new_test_names"
dct_of_pkl = {
        "C": r"\\sol.ita.chalmers.se\groups\batt_lab_data\stat_test\processed_data\Test2_1.pkl",
        "A": r"\\sol.ita.chalmers.se\groups\batt_lab_data\stat_test\processed_data\Test1_1.pkl",
        "D": r"\\sol.ita.chalmers.se\groups\batt_lab_data\stat_test\processed_data\Test2_2.pkl",
        "B": r"\\sol.ita.chalmers.se\groups\batt_lab_data\stat_test\processed_data\Test1_2.pkl"
    }
df_dict = {k: pd.read_pickle(f_) for k, f_ in dct_of_pkl.items()}
ppe_dict = {k: pp_itm(df) for k, df in df_dict.items()}
x_width = 4.5
fig_mu, axs = plt.subplots(2, 2, figsize=(2*x_width, 2*0.75*x_width), sharey='all')
[ax.yaxis.set_major_formatter(FormatStrFormatter('%.1e')) for ax in np.ravel(axs)]
plot_labels = ['(a) Test A', '(b) Test B', '(c) Test C', '(d) Test D']
for z in zip(ppe_dict.values(), axs.ravel(), plot_labels):
    z[0].plot_var_mu(z[1], sig_thresh=1.5/4)
    tmp_ax = _clean_xticks(z[1])



fig_sig, axs = plt.subplots(2, 2, figsize=(2 * x_width, 2 * 0.75 * x_width), sharey='all')
# [ax.yaxis.set_major_formatter(FormatStrFormatter('%.2f')) for ax in np.ravel(axs)]
for z in zip(ppe_dict.values(), axs.ravel(), plot_labels):
    z[0].plot_var_mu(z[1], case_to_plot='sig', sig_thresh=1.5/4)
    z[1].text(0.5, -0.165, z[2], transform=z[1].transAxes, fontsize=10, horizontalalignment='center')
    tmp_ax = _clean_xticks(z[1])
    z[1].set_xlabel('')
    z[1].set_ylabel('')
fig_sig.subplots_adjust(bottom=0.12)
fig_sig.subplots_adjust(top=0.95)
fig_sig.subplots_adjust(left=0.1)
fig_sig.subplots_adjust(hspace=0.25)
fig_sig.supxlabel('Number of Replicates')
fig_sig.supylabel('Standard error $s_n$ [%]')
if save_fig:
    fig_sig.savefig(os.path.join(base_dir, 'std_err_plot_all_tests.png'), dpi=400)

for k, pp_itm in ppe_dict.items():
    fig, ax = plt.subplots()
    ax = pp_itm.plot_var_mu(ax, case_to_plot='sig', sig_thresh=1.5 / 4)
    ax.set_ylabel('Standard error $s_n$ [%]')
    ax.set_ylim((0, 2))
    ax = _clean_xticks(ax)
    if save_fig:
        fig.tight_layout()
        fig.savefig(fr"Z:\StatisticalTest\figures_new_test_names\Test {k}\standard_error_Test_{k}.pdf")

if save_fig:
    fig_sig.tight_layout()
    fig_mu.tight_layout()
    fig_sig.savefig(r"Z:\StatisticalTest\figures\var_of_sigma.png", dpi=400)
    fig_sig.savefig(r"Z:\StatisticalTest\figures\var_of_sigma.pdf")
    fig_mu.savefig(r"Z:\StatisticalTest\figures\var_of_mu.png", dpi=400)
    fig_mu.savefig(r"Z:\StatisticalTest\figures\var_of_mu.pdf")

for k, pp_itm in ppe_dict.items():
    fig = plt.figure()
    for n, arr in pp_itm.sig_.items():
        if n < 8:
            plt.hist(arr, bins=10, density=True, alpha=0.5, label=f'{n} cells')
    plt.title(f'Test {k}')
    plt.legend()
    plt.xlabel("Variance")
    plt.ylabel("Frequency")
    if save_fig:
        fig.savefig(rf"Z:\StatisticalTest\figures\sigma_histogram_cell{k}.pdf")


lbl_fmt = 'fce_num'
# PLOT RSE FOR EACH TEST
for M, pp_it in ppe_dict.items():
    df = pp_it.rse
    if any('RPT' in c for c in df.columns):
        df.drop('RPT', axis=1, inplace=True)
    if not any('Mean' in c for c in df.index):
        df.loc['Mean', :] = df.filter(like='cells').mean()
    fig, ax = plt.subplots(figsize=(x_width, 0.75 * x_width))
    for idx in df.index:
        ax.plot(df.filter(like='cells').columns, df.loc[idx, :], label=pp_it._clean_rpt_name(idx, fmt=lbl_fmt))
    ax.set_xlabel('Number of Replicates', fontsize=12)
    ax.set_ylabel('RSE [%]', fontsize=12)
    lbls = [re.sub("[^0-9]", "", itm.get_text()) for itm in ax.get_xticklabels()]
    ax.set_xticks(ax.get_xticks())
    ax.set_xticklabels(lbls)
    ax.lines[-1].set_linestyle('dashed')
    ax.lines[-1].set_linewidth(3)
    ax.lines[-1].set_color('black')
    ax.legend(ncols=2)
    ax.grid(True)
    ax.set_ylim((0, 75))
    fig.subplots_adjust(bottom=0.14)
    if save_fig:
        fig.savefig(os.path.join(base_dir, f'Test {M}', f'rse_plot_test_{M}_{lbl_fmt}.pdf'))
        fig.savefig(os.path.join(base_dir, f'Test {M}', f'rse_plot_test_{M}_{lbl_fmt}.png'), dpi=400)


 # OUTPUT RSE TO TABLE IN POWERPOINT AND LATEX FORMAT
output_ppt_table = 1
if output_ppt_table:
    for M in ['A', 'B', 'C', 'D']:
        df = ppe_dict[M].rse
        df.fillna(0, inplace=True)
        df.columns = [c.replace('_test', '') for c in df.columns]
        df.loc['Mean', :] = df.filter(like='cells').mean()
        with open(os.path.join(base_dir, f'Test {M}', 'rse_tab.tex'), 'w') as tf:
            tf.write(df.to_latex(float_format="%.0f"))
        # df.insert(0, 'RPT', df.index)
        m, n = df.shape
        col_format = ['.2' for k in range(n)]
        df_to_powerpoint(os.path.join(base_dir, f'Test {M}', f"rse_table_test_{M}.pptx"),
                         df.iloc[:, :-1],
                         col_formatters=col_format,
                         width=19,
                         height=6,
                         left=3,
                         top=3,
                         name='risk_table')
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pr_shp = df_to_table(slide, df, col_formatters=col_format)

